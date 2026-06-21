"""PPTX export service — renders presentation JSON into PowerPoint files."""

import json
import os
import logging
from datetime import datetime, timezone
from typing import Optional

from pptx import Presentation as PptxPresentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE

from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select

from app.models.models import Export, Presentation, PresentationStatus
from app.core.config import settings

logger = logging.getLogger(__name__)

# ── Color themes ────────────────────────────────

THEMES = {
    "business": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0x1A, 0x56, 0xDB),  # Blue
        "dark": RGBColor(0x1E, 0x1E, 0x2E),
        "text": RGBColor(0x33, 0x33, 0x33),
        "light_bg": RGBColor(0xF5, 0xF7, 0xFA),
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "dark_invest": {
        "bg": RGBColor(0x1A, 0x1A, 0x2E),
        "accent": RGBColor(0x00, 0xD2, 0x8E),  # Green
        "dark": RGBColor(0x0D, 0x0D, 0x1A),
        "text": RGBColor(0xE0, 0xE0, 0xE0),
        "light_bg": RGBColor(0x24, 0x24, 0x3E),
        "title_color": RGBColor(0x00, 0xD2, 0x8E),
    },
    "modern": {
        "bg": RGBColor(0xF8, 0xF9, 0xFA),
        "accent": RGBColor(0x6C, 0x5C, 0xE7),  # Purple
        "dark": RGBColor(0x2D, 0x2D, 0x3F),
        "text": RGBColor(0x33, 0x33, 0x33),
        "light_bg": RGBColor(0xEE, 0xF0, 0xF4),
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "marketing": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0xE6, 0x4A, 0x19),  # Orange
        "dark": RGBColor(0x2D, 0x2D, 0x2D),
        "text": RGBColor(0x33, 0x33, 0x33),
        "light_bg": RGBColor(0xFF, 0xF4, 0xEB),
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
    },
    "analytical": {
        "bg": RGBColor(0xFF, 0xFF, 0xFF),
        "accent": RGBColor(0x00, 0x96, 0x88),  # Teal
        "dark": RGBColor(0x26, 0x32, 0x38),
        "text": RGBColor(0x33, 0x33, 0x33),
        "light_bg": RGBColor(0xE0, 0xF2, 0xF1),
        "title_color": RGBColor(0xFF, 0xFF, 0xFF),
    },
}

DEFAULT_THEME = THEMES["business"]

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


class PptxRenderer:
    """Renders presentation JSON into a PPTX file."""

    def __init__(self, theme_name: str = "business"):
        self.theme = THEMES.get(theme_name, DEFAULT_THEME)

    async def render(
        self,
        presentation_json: str,
        output_path: str,
    ) -> str:
        """Render presentation JSON to PPTX file. Returns output path."""
        data = json.loads(presentation_json)
        prs = PptxPresentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        slides = data.get("slides", [])
        for slide_data in slides:
            self._add_slide(prs, slide_data)

        os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
        prs.save(output_path)
        return output_path

    def _add_slide(self, prs: PptxPresentation, slide_data: dict):
        """Add a single slide."""
        slide_type = slide_data.get("type", "thesis")
        title = slide_data.get("title", "")
        content = slide_data.get("content", [])
        subtitle = slide_data.get("subtitle", "")

        layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(layout)

        if slide_type == "cover":
            self._render_cover(slide, title, subtitle)
        elif slide_type in ("conclusions", "cta"):
            self._render_conclusions(slide, title, content)
        else:
            self._render_content_slide(slide, title, content, slide_type)

    def _render_cover(self, slide, title: str, subtitle: str):
        """Render a cover slide."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self.theme["dark"]

        # Title
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(44)
        p.font.color.rgb = self.theme["title_color"]
        p.font.bold = True

        if subtitle:
            txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(1))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = subtitle
            p2.font.size = Pt(20)
            p2.font.color.rgb = RGBColor(0xAA, 0xAA, 0xBB)

    def _render_content_slide(self, slide, title: str, content: list, slide_type: str):
        """Render a standard content slide."""
        # Accent bar
        bar = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(0.15), SLIDE_HEIGHT
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.theme["accent"]
        bar.line.fill.background()

        # Title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.7))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(28)
        p.font.color.rgb = self.theme["dark"]
        p.font.bold = True

        # Content bullets
        if isinstance(content, list) and content:
            txBox2 = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11), Inches(5))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for i, bullet in enumerate(content):
                if i == 0:
                    p2 = tf2.paragraphs[0]
                else:
                    p2 = tf2.add_paragraph()
                p2.text = str(bullet)[:200]
                p2.font.size = Pt(16)
                p2.font.color.rgb = self.theme["text"]
                p2.space_after = Pt(8)
                p2.level = 0

    def _render_conclusions(self, slide, title: str, content: list):
        """Render conclusions/CTA slide with accent background."""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self.theme["accent"]

        # Title
        txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(1))
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Content
        if isinstance(content, list) and content:
            txBox2 = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(9), Inches(3))
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            for i, bullet in enumerate(content):
                if i == 0:
                    p2 = tf2.paragraphs[0]
                else:
                    p2 = tf2.add_paragraph()
                p2.text = str(bullet)
                p2.font.size = Pt(18)
                p2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                p2.alignment = PP_ALIGN.CENTER
                p2.space_after = Pt(6)


class ExportService:
    """Orchestrates export of presentations to PPTX/PDF."""

    def __init__(self):
        self.renderer = PptxRenderer()

    async def export_pptx(
        self,
        presentation_id: int,
        user_id: int,
        db: AsyncSession,
        theme: str = "business",
    ) -> dict:
        """Export presentation to PPTX and save."""
        result = await db.execute(
            select(Presentation).where(Presentation.id == presentation_id, Presentation.user_id == user_id)
        )
        pres = result.scalar_one_or_none()
        if not pres:
            raise ValueError("Presentation not found")
        if not pres.presentation_json:
            raise ValueError("Presentation has no slides data")

        # Create export record
        export = Export(
            presentation_id=presentation_id,
            user_id=user_id,
            format="pptx",
            status="rendering",
        )
        db.add(export)
        await db.commit()
        await db.refresh(export)

        try:
            # Render PPTX
            ext = "pptx"
            filename = f"{pres.title.replace(' ', '_')[:50]}_{export.id}.{ext}"
            output_dir = os.path.join(settings.upload_dir, "exports", str(user_id))
            output_path = os.path.join(output_dir, filename)

            renderer = PptxRenderer(theme)
            output_path = await renderer.render(pres.presentation_json, output_path)

            # Update export
            export.file_path = output_path
            export.status = "completed"
            await db.commit()

            return {
                "id": export.id,
                "format": "pptx",
                "status": "completed",
                "file_path": output_path,
                "filename": filename,
            }

        except Exception as e:
            export.status = "failed"
            await db.commit()
            logger.error(f"Export failed for {presentation_id}: {e}", exc_info=True)
            raise

    async def export_pdf(
        self,
        presentation_id: int,
        user_id: int,
        db: AsyncSession,
        theme: str = "business",
    ) -> dict:
        """Export presentation to PDF (convert from PPTX using LibreOffice)."""
        # First generate PPTX
        pptx_result = await self.export_pptx(presentation_id, user_id, db, theme)

        # Try to convert to PDF using LibreOffice
        pptx_path = pptx_result["file_path"]
        pdf_path = pptx_path.replace(".pptx", ".pdf")

        try:
            import subprocess
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir",
                 os.path.dirname(pdf_path), pptx_path],
                capture_output=True, text=True, timeout=60
            )
            if result.returncode == 0 and os.path.exists(pdf_path):
                # Create PDF export record
                export = Export(
                    presentation_id=presentation_id,
                    user_id=user_id,
                    format="pdf",
                    status="completed",
                    file_path=pdf_path,
                )
                db.add(export)
                await db.commit()
                await db.refresh(export)

                return {
                    "id": export.id,
                    "format": "pdf",
                    "status": "completed",
                    "file_path": pdf_path,
                    "filename": os.path.basename(pdf_path),
                }
            else:
                logger.warning(f"LibreOffice conversion failed: {result.stderr[:200]}")
                raise ValueError("PDF conversion failed: LibreOffice not available or error")
        except FileNotFoundError:
            logger.warning("LibreOffice not installed, cannot convert to PDF")
            raise ValueError("PDF conversion requires LibreOffice")
        except Exception as e:
            logger.error(f"PDF export failed: {e}")
            raise


export_service = ExportService()
